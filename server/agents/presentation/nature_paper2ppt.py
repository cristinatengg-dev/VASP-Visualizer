#!/usr/bin/env python3
import json
import os
import re
import sys
from datetime import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
NAVY = RGBColor(10, 17, 40)
GREEN = RGBColor(5, 150, 105)
BLUE = RGBColor(37, 99, 235)
GRAY = RGBColor(100, 116, 139)
LIGHT = RGBColor(248, 250, 252)


def clean(text, fallback=""):
    if text is None:
        return fallback
    return re.sub(r"\s+", " ", str(text)).strip() or fallback


def truncate(text, limit=150):
    text = clean(text)
    return text if len(text) <= limit else text[: limit - 1] + "..."


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=NAVY, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align:
      p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.6, 0.35, 8.9, 0.45, title, size=22, bold=True, color=NAVY)
    if subtitle:
        add_textbox(slide, 0.62, 0.88, 11.4, 0.35, subtitle, size=10, color=GRAY)
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.24), Inches(12.1), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(226, 232, 240)
    line.line.color.rgb = RGBColor(226, 232, 240)


def add_bullets(slide, x, y, w, h, items, size=13, color=NAVY):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = truncate(item, 210)
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)
        p._p.get_or_add_pPr().insert(0, p._element.makeelement("{http://schemas.openxmlformats.org/drawingml/2006/main}buChar", {"char": "•"}))
    return box


def add_metric(slide, x, y, label, value, color=BLUE):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.45), Inches(0.86))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = RGBColor(226, 232, 240)
    add_textbox(slide, x + 0.16, y + 0.12, 2.1, 0.2, label, size=8, bold=True, color=GRAY)
    add_textbox(slide, x + 0.16, y + 0.38, 2.1, 0.3, value, size=15, bold=True, color=color)


def first_items(items, limit):
    return list(items or [])[:limit]


def generate(payload, out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    prompt = clean(payload.get("prompt"), "科研计算工作流")
    research = payload.get("research") or {}
    research_stack = payload.get("researchStack") or {}
    papers = research.get("papers") or []
    idea = payload.get("selectedIdea") or {}
    blueprint = idea.get("blueprint") or {}
    structure_source = blueprint.get("structure_source") or {}
    model = payload.get("modelStructure") or {}
    compute_intent = payload.get("computeIntent") or {}
    compiled_files = payload.get("compiledFiles") or {}
    job_status = payload.get("jobStatus") or {}
    compute_result = payload.get("computeResult") or {}

    # 1 Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 0.65, 0.8, 8.7, 0.7, "EliangMat AI Research Report", size=30, bold=True, color=NAVY)
    add_textbox(slide, 0.68, 1.55, 10.8, 0.7, truncate(prompt, 150), size=16, color=GRAY)
    add_metric(slide, 0.7, 3.0, "Papers", str(len(papers)), GREEN)
    add_metric(slide, 3.35, 3.0, "Model atoms", str(model.get("atomCount") or 0), BLUE)
    add_metric(slide, 6.0, 3.0, "Input files", str(len(compiled_files)), RGBColor(217, 119, 6))
    add_metric(slide, 8.65, 3.0, "Job", clean(job_status.get("status"), "-"), RGBColor(190, 18, 60))
    if research_stack.get("feasibility"):
        add_textbox(slide, 11.25, 3.12, 1.2, 0.18, "Feas.", size=8, bold=True, color=GRAY)
        add_textbox(slide, 11.25, 3.42, 1.2, 0.3, str(research_stack.get("feasibility", {}).get("score", "-")), size=15, bold=True, color=NAVY)
    add_textbox(slide, 0.7, 6.65, 6.5, 0.3, f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M')}", size=9, color=GRAY)
    add_textbox(slide, 8.8, 6.55, 3.9, 0.35, "Evidence-backed PPT workflow", size=9, color=GRAY, align=PP_ALIGN.RIGHT)

    # 2 Workflow thinking summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Visible Reasoning Summary", "From prompt to evidence-backed modeling and compute decisions")
    add_bullets(slide, 0.8, 1.65, 11.7, 4.8, [
        "Interpret the task as a literature-grounded catalyst search plus atomistic model construction workflow.",
        "Search literature sources first, then rank model candidates by evidence, structure availability, and compute feasibility.",
        "Ask for model confirmation before building, because unsupported material guesses can corrupt downstream DFT inputs.",
        "Compile engine-specific input files only after software selection, then require approval before job submission.",
        "Generate the PPT only after compute completion or explicit user request, preserving the full evidence chain.",
    ], size=15)

    # 3 Literature evidence
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Literature Evidence", "Top retrieved records used by the recommendation step")
    if papers:
        add_bullets(slide, 0.75, 1.55, 11.8, 5.3, [
            f"{clean(p.get('title'), 'Untitled')} ({clean(p.get('year'), 'n.d.')}) · {clean(p.get('source'), 'source')} · {clean(p.get('doi'), clean(p.get('url'), 'no DOI'))}"
            for p in first_items(papers, 7)
        ], size=11)
    else:
        add_bullets(slide, 0.75, 1.55, 11.8, 5.3, ["No paper evidence was available in the current run."], size=13)

    # 4 Synthesis feasibility and experiment design
    if research_stack:
        synthesis = research_stack.get("synthesis") or {}
        feasibility = research_stack.get("feasibility") or {}
        experiment = research_stack.get("experiment") or {}
        routes = synthesis.get("routes") or []
        route = routes[0] if routes else {}
        first_batch = experiment.get("first_batch") or []
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, "Synthesis Feasibility", "Routes, feasibility score, and first-batch experiment design")
        add_metric(slide, 0.75, 1.55, "Score", str(feasibility.get("score", "-")), GREEN)
        add_metric(slide, 3.4, 1.55, "Level", clean(feasibility.get("level"), "-"), BLUE)
        add_metric(slide, 6.05, 1.55, "Routes", str(len(routes)), RGBColor(217, 119, 6))
        add_metric(slide, 8.7, 1.55, "Experiments", str(len(first_batch)), RGBColor(190, 18, 60))
        add_textbox(slide, 0.8, 2.78, 5.55, 0.28, "Synthesis route", size=13, bold=True, color=NAVY)
        add_bullets(slide, 0.8, 3.16, 5.65, 2.95, [
            f"Route: {clean(route.get('title'), 'No route')}",
            f"Method: {clean(route.get('method'), 'n/a')}",
            f"Source: {clean(route.get('source'), 'heuristic')} {clean(route.get('doi'), '')}".strip(),
            f"Precursors: {clean(', '.join(route.get('precursors') or []), 'n/a')}",
            f"Conditions: {clean((route.get('conditions') or {}).get('temperature'), 'n/a')}; {clean((route.get('conditions') or {}).get('atmosphere'), 'n/a')}",
            f"Risk: {clean(route.get('risk'), 'n/a')}",
        ], size=10)
        add_textbox(slide, 7.05, 2.78, 5.4, 0.28, "Experiment design", size=13, bold=True, color=BLUE)
        experiment_lines = []
        for row in first_items(first_batch, 5):
            experiment_lines.append(", ".join([f"{k}={v}" for k, v in list(row.items())[:5]]))
        add_bullets(slide, 7.05, 3.16, 5.35, 2.95, experiment_lines or ["No experiment matrix recorded"], size=10)

    # 5 Model recommendation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Model Recommendation", "Recommended starter structure and why it was selected")
    add_textbox(slide, 0.75, 1.55, 5.3, 0.4, clean(idea.get("title"), "Recommended model"), size=20, bold=True, color=NAVY)
    add_bullets(slide, 0.75, 2.1, 5.6, 3.7, [
        f"Formula/source: {clean(structure_source.get('formula'), clean(idea.get('material_family'), 'n/a'))}",
        f"Model type: {clean(idea.get('recommended_model_type'), 'starter')}",
        f"Reason: {clean(idea.get('fit_reason'), clean(blueprint.get('why_this_idea'), 'n/a'))}",
        f"Caution: {clean('; '.join(blueprint.get('caution_notes') or []), 'Validate against target papers before publication use')}",
    ], size=12)
    add_textbox(slide, 7.0, 1.7, 5.3, 0.35, "What can be calculated", size=13, bold=True, color=BLUE)
    add_bullets(slide, 7.0, 2.15, 5.2, 3.6, [
        clean(blueprint.get("what_can_be_calculated"), "Relaxation, adsorption energy, charge transfer, DOS or band structure as appropriate."),
        f"First step: {clean(blueprint.get('first_step'), 'Build and relax the starter structure')}",
        f"Second step: {clean(blueprint.get('second_step'), 'Run validation or property calculation')}",
    ], size=12)

    # 5 Built structure
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Deterministic Structure Build", "Modeling output passed to compute")
    add_metric(slide, 0.75, 1.65, "Filename", clean(model.get("filename"), "structure"), BLUE)
    add_metric(slide, 3.4, 1.65, "Atoms", str(model.get("atomCount") or 0), GREEN)
    add_metric(slide, 6.05, 1.65, "Bonds", str(model.get("bondCount") or 0), RGBColor(217, 119, 6))
    add_metric(slide, 8.7, 1.65, "Periodic", "yes" if model.get("latticeVectors") else "no", RGBColor(190, 18, 60))
    add_bullets(slide, 0.85, 3.0, 11.5, 3.2, [
        "The structure was generated before image or compute steps, reducing chemistry hallucination risk.",
        "Downstream inputs are compiled from atom coordinates rather than from visual prompts.",
        "If the model is not chemically appropriate, revise the model stage before using the compute package.",
    ], size=14)

    # 6 Compute inputs
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Compute Input Package", "Software selection and generated input files")
    add_bullets(slide, 0.75, 1.5, 5.5, 4.9, [
        f"Engine: {clean(compute_intent.get('engine'), 'vasp')}",
        f"Workflow: {clean(compute_intent.get('workflow'), 'relax')}",
        f"Quality: {clean(compute_intent.get('quality'), 'standard')}",
        f"Spin mode: {clean(compute_intent.get('spin_mode'), 'auto')}",
        f"vdW: {compute_intent.get('vdw', True)}",
    ], size=13)
    add_textbox(slide, 7.0, 1.55, 4.9, 0.35, "Generated files", size=13, bold=True, color=BLUE)
    add_bullets(slide, 7.0, 2.05, 5.0, 3.7, list(compiled_files.keys())[:12] or ["No files recorded"], size=12)

    # 7 Compute result
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Compute Result", "Job status and extracted metrics")
    add_metric(slide, 0.75, 1.6, "Status", clean(job_status.get("status"), "-"), GREEN)
    add_metric(slide, 3.4, 1.6, "Energy (eV)", clean(compute_result.get("totalEnergyEv"), "N/A"), BLUE)
    add_metric(slide, 6.05, 1.6, "Converged", clean(compute_result.get("converged"), "N/A"), RGBColor(217, 119, 6))
    add_metric(slide, 8.7, 1.6, "Max force", clean(compute_result.get("maxForceEvPerA"), "N/A"), RGBColor(190, 18, 60))
    add_bullets(slide, 0.85, 3.0, 11.5, 3.2, [
        f"Profile: {clean(job_status.get('profileId'), 'n/a')}",
        f"External job id: {clean(job_status.get('externalJobId'), 'local/demo')}",
        "Use this result as a workflow checkpoint; publication-grade conclusions require convergence and methodological validation.",
    ], size=14)

    # 8 Next actions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Recommended Next Actions", "Follow-up work before manuscript or figure production")
    add_bullets(slide, 0.8, 1.55, 11.6, 5.1, [
        "Confirm that the selected structure matches the target catalyst chemistry and relevant experimental phase.",
        "Run a high-quality relaxation after the starter calculation if the first job used a fast/demo profile.",
        "For CO2 hydrogenation, compare CO2, H2, formate, COOH, CO and product-state adsorption energies on the same surface model.",
        "Use charge density difference, DOS/PDOS, Bader charge, or reaction pathway calculations only after the structure is validated.",
        "Keep article DOIs and structure source IDs attached to every figure caption.",
    ], size=14)

    # 9 Appendix
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Appendix: Evidence Links", "Sources retained from the agent run")
    links = []
    for p in first_items(papers, 10):
        if p.get("doi"):
            links.append(f"DOI: {p.get('doi')}")
        elif p.get("url"):
            links.append(f"URL: {p.get('url')}")
    add_bullets(slide, 0.8, 1.55, 11.6, 5.1, links or ["No external evidence links were retained."], size=11)

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    return {
        "slide_count": len(prs.slides),
        "qa": f"{len(prs.slides)} slides generated with paper/model/compute evidence sections.",
    }


def main():
    if len(sys.argv) != 2:
        raise SystemExit("Usage: nature_paper2ppt.py <out.pptx>")
    payload = json.load(sys.stdin)
    result = generate(payload, sys.argv[1])
    print(json.dumps(result, ensure_ascii=False))


if __name__ == "__main__":
    main()
